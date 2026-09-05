from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
GREEN_DARK = RGBColor(26, 71, 42)
GREEN_LIGHT = RGBColor(67, 160, 71)
GREEN_ACCENT = RGBColor(46, 125, 50)
GOLD = RGBColor(255, 193, 7)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(44, 62, 80)

def add_image_from_url(slide, url, left, top, width, height):
    try:
        response = requests.get(url, timeout=5)
        if response.status_code == 200:
            img = BytesIO(response.content)
            slide.shapes.add_picture(img, left, top, width=width, height=height)
            return True
    except:
        pass
    return False

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_DARK
    
    # Add decorative shapes
    shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(4), Inches(4))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = GREEN_LIGHT
    shape1.fill.transparency = 0.3
    shape1.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GOLD
    shape2.fill.transparency = 0.2
    shape2.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = GOLD
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Helping Farmers, Protecting Futures"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.italic = True
    tagline_para.font.color.rgb = GREEN_LIGHT
    tagline_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items, farmer_image_url=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250)
    
    # Title bar with gradient effect
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = GREEN_ACCENT
    title_shape.line.color.rgb = GREEN_ACCENT
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(7), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Add image if provided
    if farmer_image_url:
        add_image_from_url(slide, farmer_image_url, Inches(5.5), Inches(1.2), Inches(4), Inches(3.5))
        content_width = 5
    else:
        content_width = 8.4
    
    # Content
    y_pos = 1.3
    for item in content_items:
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(content_width), Inches(1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_TEXT
        para.space_before = Pt(6)
        para.space_after = Pt(6)
        y_pos += 0.95
    
    return slide

# Slide 1: Title
add_title_slide("🌾 Kisan-Setu AI", "Smart Crop Insurance Portal")

# Slide 2: What is Kisan-Setu AI
add_content_slide("What is Kisan-Setu AI?", [
    "✓ AI-powered crop damage assessment system",
    "✓ Automated insurance claim verification",
    "✓ Real-time satellite NDVI analysis",
    "✓ Geo-tagged farm verification",
    "✓ Instant claim processing"
])

# Slide 3: Features
add_content_slide("Key Features", [
    "📸 Live Farm Photo Capture - Verify your crop damage",
    "🛰️ Satellite Damage Analysis - AI-powered NDVI scoring",
    "⚡ Instant Results - Claim approval within seconds",
    "💰 Direct Payout - Insurance amount credited instantly",
    "📱 Mobile Friendly - Easy access for farmers"
])

# Slide 4: How It Works
add_content_slide("Claim Process", [
    "1️⃣ Enter Aadhaar/PAN number",
    "2️⃣ Upload clear photo of your damaged crops",
    "3️⃣ AI verifies Bhulekh & geolocation",
    "4️⃣ Satellite NDVI analysis for damage %",
    "5️⃣ Instant approval/rejection with reasons"
])

# Slide 5: Benefits
add_content_slide("Benefits", [
    "💚 Reduces corruption - Fully automated verification",
    "⏱️ Saves time - No paperwork required",
    "🤝 Fair assessment - AI-powered decisions",
    "📊 Transparent - Detailed reason explanations",
    "🚀 Fast payout - Direct bank transfer"
])

# Slide 6: Coverage Areas
add_content_slide("Districts Covered", [
    "🏛️ Lucknow District - Multiple Tehsils",
    "   • Mohanlalganj Tehsil",
    "   • Sadar (Lucknow) Tehsil",
    "   • Malihabad Tehsil",
    "📈 Coverage expanding to more districts soon"
])

# Slide 7: Supported Crops
add_content_slide("Supported Crops", [
    "🌾 Wheat & Rice",
    "🌽 Maize & Pulses",
    "🥬 Vegetables & Fruits",
    "🌻 Oilseeds & Cash Crops",
    "✅ More crops being added regularly"
])

# Slide 8: Documentation Required
add_content_slide("Documents Required", [
    "✓ Aadhaar Card / PAN Card",
    "✓ Clear photo of damaged farm (geotag enabled)",
    "✓ Khasra/Khata number from Bhulekh",
    "✓ Bank account details (for payout)",
    "✓ Contact number for verification"
])

# Slide 9: Rejection Reasons
add_content_slide("Claim Rejection Reasons", [
    "❌ Photo location mismatch (outside farm boundary)",
    "❌ Aadhaar vs Bhulekh name mismatch",
    "❌ Damage below threshold (<30%)",
    "❌ Blurry/unclear photos",
    "❌ Already claimed for this season"
])

# Slide 10: Success Stories
add_content_slide("Our Impact", [
    "📊 128 Total Claims Processed",
    "✅ 102 Claims Approved (79.7% success rate)",
    "❌ 26 Claims Rejected (with reasons)",
    "💰 ₹15,25,000 Total Payout",
    "😊 10,000+ Happy Farmers"
])

# Slide 11: How to Get Started
add_content_slide("Get Started Now!", [
    "1. Click 'Proceed to Claim Form' below",
    "2. Enter your Aadhaar number",
    "3. Upload a clear farm photo",
    "4. Get instant results",
    "5. Receive payout in your bank account"
])

# Slide 12: Contact & Support
add_content_slide("Support & Contact", [
    "📞 Helpline: 1800-KISAN-SETU",
    "📧 Email: support@kisansetu.gov.in",
    "🌐 Website: www.kisansetu.gov.in",
    "💬 WhatsApp Support: Available 24/7",
    "🏢 Visit: Your nearest Tehsil Office"
])

# Save presentation
prs.save('presentation.pptx')
print("✅ Presentation created successfully: presentation.pptx")
